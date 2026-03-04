# =============================================================================
# modes/slides_mode.py — Слайды (JPG/PNG + PDF + PPTX)
# =============================================================================
"""
SlidesMode v2:
  - Загружает изображения (.jpg/.png/.bmp/.webp)
  - Загружает PDF (через pymupdf/fitz) — каждая страница = слайд
  - Загружает PPTX (через python-pptx + LibreOffice или COM-конвертация)
  - Свайп открытой ладони вправо/влево → навигация
  - Плавная анимация перехода (ease-out cubic)
"""

import os
import cv2
import numpy as np
import time
from typing import Optional

from modes.base_mode import BaseMode
from modules.hand_tracker import HandData
from modules.gesture_detector import Gesture
from modules.ui import UIRenderer, put_text_centered, draw_rounded_rect
from config import (
    MODE_MENU, SLIDES_DIR,
    COLOR_WHITE, COLOR_ACCENT, COLOR_GRAY,
    HUD_HEIGHT, FONT_MAIN, FONT_BOLD, FONT_SMALL
)

SUPPORTED_IMG_EXTS = ('.jpg', '.jpeg', '.png', '.bmp', '.webp', '.tiff')
ANIM_DURATION      = 0.35

# --- Опциональный импорт PDF-движка (pymupdf) ---
try:
    import fitz  # pymupdf
    _FITZ_OK = True
except ImportError:
    _FITZ_OK = False

# --- Опциональный импорт PPTX ---
try:
    from pptx import Presentation
    from pptx.util import Inches
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


def _load_pdf(path: str, dpi: int = 150) -> list[np.ndarray]:
    """Конвертирует каждую страницу PDF в BGR-изображение."""
    if not _FITZ_OK:
        print("[Slides] pymupdf не установлен — PDF не поддерживается.")
        return []
    pages = []
    try:
        doc = fitz.open(path)
        scale = dpi / 72.0
        mat   = fitz.Matrix(scale, scale)
        for page in doc:
            pix  = page.get_pixmap(matrix=mat)
            img  = np.frombuffer(pix.samples, dtype=np.uint8)
            img  = img.reshape(pix.height, pix.width, pix.n)
            if pix.n == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2BGR)
            else:
                img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)
            pages.append(img)
        doc.close()
    except Exception as e:
        print(f"[Slides] Ошибка чтения PDF {path}: {e}")
    return pages


def _load_pptx(path: str) -> list[np.ndarray]:
    """
    Попытка конвертировать PPTX в изображения.
    Стратегия 1: LibreOffice CLI → временный PDF → fitz
    Стратегия 2: Если LibreOffice нет — python-pptx заглушка с именами слайдов
    """
    import subprocess, tempfile, shutil

    # --- Стратегия 1: LibreOffice ---
    libre = shutil.which("soffice") or shutil.which("libreoffice")
    if libre and _FITZ_OK:
        tmpdir = tempfile.mkdtemp()
        try:
            result = subprocess.run(
                [libre, "--headless", "--convert-to", "pdf",
                 "--outdir", tmpdir, path],
                capture_output=True, timeout=60
            )
            pdf_name = os.path.splitext(os.path.basename(path))[0] + ".pdf"
            pdf_path = os.path.join(tmpdir, pdf_name)
            if os.path.isfile(pdf_path):
                slides = _load_pdf(pdf_path)
                print(f"[Slides] PPTX→PDF через LibreOffice: {len(slides)} слайдов")
                return slides
        except Exception as e:
            print(f"[Slides] LibreOffice конвертация не удалась: {e}")
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    # --- Стратегия 2: python-pptx + заглушка ---
    if _PPTX_OK:
        try:
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides):
                # Создаём заглушку с именем слайда
                img = np.full((720, 1280, 3), 20, dtype=np.uint8)
                title_shape = None
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        title_shape = shape
                        break
                title = title_shape.text[:60] if title_shape else f"Слайд {i+1}"
                cv2.putText(img, f"[PPTX] {title}", (50, 360),
                            cv2.FONT_HERSHEY_SIMPLEX, 1.0, (0, 200, 255), 2, cv2.LINE_AA)
                cv2.putText(img, "Установите LibreOffice для полного рендера",
                            (50, 410), cv2.FONT_HERSHEY_PLAIN, 1.5, (80,80,80), 1, cv2.LINE_AA)
                slides.append(img)
            print(f"[Slides] PPTX-заглушка: {len(slides)} слайдов (нет рендера)")
            return slides
        except Exception as e:
            print(f"[Slides] Ошибка чтения PPTX {path}: {e}")

    print("[Slides] PPTX: установите python-pptx или LibreOffice.")
    return []


class SlidesMode(BaseMode):
    """Режим слайдов: изображения + PDF + PPTX."""

    def __init__(self, ui: UIRenderer):
        super().__init__(ui)
        self._slides: list[np.ndarray] = []
        self._filenames: list[str]     = []
        self._index       = 0
        self._anim_offset = 0.0
        self._anim_dir    = 0
        self._anim_start  = 0.0

    # ------------------------------------------------------------------
    # Lifecycle
    # ------------------------------------------------------------------

    def on_enter(self):
        self._load_slides()
        self._index       = 0
        self._anim_offset = 0.0
        self._anim_dir    = 0

    def on_exit(self):
        pass

    # ------------------------------------------------------------------
    # Загрузка слайдов
    # ------------------------------------------------------------------

    def _load_slides(self):
        self._slides    = []
        self._filenames = []

        if not os.path.isdir(SLIDES_DIR):
            os.makedirs(SLIDES_DIR, exist_ok=True)
            return

        files = sorted(os.listdir(SLIDES_DIR))

        for fname in files:
            path = os.path.join(SLIDES_DIR, fname)
            ext  = os.path.splitext(fname)[1].lower()

            if ext in SUPPORTED_IMG_EXTS:
                img = cv2.imread(path)
                if img is not None:
                    self._slides.append(img)
                    self._filenames.append(fname)

            elif ext == '.pdf':
                pages = _load_pdf(path)
                for j, page in enumerate(pages):
                    self._slides.append(page)
                    self._filenames.append(f"{fname} [{j+1}]")

            elif ext in ('.pptx', '.ppt'):
                pages = _load_pptx(path)
                for j, page in enumerate(pages):
                    self._slides.append(page)
                    self._filenames.append(f"{fname} [{j+1}]")

        print(f"[Слайды] Загружено {len(self._slides)} слайдов из {SLIDES_DIR}")

    # ------------------------------------------------------------------
    # Update
    # ------------------------------------------------------------------

    def update(
        self,
        frame: np.ndarray,
        hands: list[HandData],
        gesture: str,
        primary_hand: Optional[HandData],
    ) -> Optional[str]:
        fh, fw = frame.shape[:2]

        back = self._check_back_gesture(gesture)
        if back:
            return back

        # Навигация свайпом
        if self._anim_start == 0.0:
            if gesture == Gesture.SWIPE_RIGHT:
                self._navigate(+1)
            elif gesture == Gesture.SWIPE_LEFT:
                self._navigate(-1)

        if not self._slides:
            self._draw_empty_state(frame, fw, fh)
        else:
            self._draw_slides(frame, fw, fh)

        self._ui.draw_hud(frame, "Слайды", 0, "", len(hands))
        self._ui.draw_back_hint(frame)
        self._draw_nav_hint(frame, fw, fh)

        if primary_hand:
            ix, iy = primary_hand.index_tip_px
            self._ui.draw_finger_cursor(frame, ix, iy)

        return None

    # ------------------------------------------------------------------
    # Навигация
    # ------------------------------------------------------------------

    def _navigate(self, direction: int):
        n = len(self._slides)
        if n == 0:
            return
        new_idx = self._index + direction
        # КЛAMP — не оборачиваем, стоп на границах
        new_idx = max(0, min(n - 1, new_idx))
        if new_idx == self._index:
            return  # Уже на границе — игнорируем
        self._anim_dir   = direction
        self._anim_start = time.time()
        self._index      = new_idx

    # ------------------------------------------------------------------
    # Отрисовка
    # ------------------------------------------------------------------

    def _draw_slides(self, frame: np.ndarray, fw: int, fh: int):
        available_h = fh - HUD_HEIGHT - 80
        available_w = fw

        now      = time.time()
        elapsed  = now - self._anim_start if self._anim_start else 0
        progress = min(elapsed / ANIM_DURATION, 1.0) if self._anim_start else 1.0
        t        = 1.0 - (1.0 - progress) ** 3   # ease-out cubic

        slide_offset = int((1.0 - t) * fw * self._anim_dir * -1)
        if progress >= 1.0:
            self._anim_start = 0.0
            self._anim_dir   = 0
            slide_offset     = 0

        cur_img  = self._fit_slide(self._slides[self._index], available_w, available_h)
        ch, cw   = cur_img.shape[:2]
        cx       = (fw - cw) // 2 + slide_offset
        cy       = HUD_HEIGHT + (available_h - ch) // 2

        self._blit_safe(frame, cur_img, cx, cy)
        draw_rounded_rect(frame, cx-2, cy-2, cx+cw+2, cy+ch+2, 4, COLOR_ACCENT, thickness=1)

        n = len(self._slides)
        put_text_centered(frame, f"{self._index + 1} / {n}",
                          fw // 2, fh - HUD_HEIGHT + 10, FONT_BOLD, 0.9, COLOR_WHITE, 1)
        fname = self._filenames[self._index]
        put_text_centered(frame, fname, fw // 2, fh - HUD_HEIGHT + 36,
                          FONT_SMALL, 1.1, COLOR_GRAY, 1)
        self._draw_dots(frame, fw, fh, n)

    def _fit_slide(self, img: np.ndarray, max_w: int, max_h: int) -> np.ndarray:
        h, w  = img.shape[:2]
        scale = min(max_w / w, max_h / h, 1.0)
        if scale < 1.0:
            return cv2.resize(img, (int(w*scale), int(h*scale)), interpolation=cv2.INTER_AREA)
        return img

    @staticmethod
    def _blit_safe(dst: np.ndarray, src: np.ndarray, x: int, y: int):
        dh, dw = dst.shape[:2]
        sh, sw = src.shape[:2]
        sx1 = max(0, -x);  sy1 = max(0, -y)
        sx2 = min(sw, dw-x); sy2 = min(sh, dh-y)
        dx1 = max(0, x);   dy1 = max(0, y)
        dx2 = dx1+(sx2-sx1); dy2 = dy1+(sy2-sy1)
        if sx2 > sx1 and sy2 > sy1:
            dst[dy1:dy2, dx1:dx2] = src[sy1:sy2, sx1:sx2]

    def _draw_dots(self, frame: np.ndarray, fw: int, fh: int, n: int):
        MAX_DOTS = 20
        display_n = min(n, MAX_DOTS)
        dot_r = 5; gap = 14
        total  = display_n * (dot_r*2) + (display_n-1) * gap
        start_x = (fw - total) // 2
        cy = fh - 48
        for i in range(display_n):
            # Если слайдов много — показываем относительный индикатор
            real_i = int(i * n / display_n)
            color  = COLOR_ACCENT if real_i == self._index or (
                        i == display_n - 1 and self._index >= int((display_n-1)*n/display_n)
                     ) else COLOR_GRAY
            cx = start_x + i * (dot_r*2 + gap) + dot_r
            cv2.circle(frame, (cx, cy), dot_r, color, -1, cv2.LINE_AA)

    def _draw_nav_hint(self, frame: np.ndarray, fw: int, fh: int):
        pdf_hint = ""
        if not _FITZ_OK:
            pdf_hint = "  [PDF: pip install pymupdf]"
        hint = f"< Vzмах vlevo  /  Vzмakh vpravo >{pdf_hint}"
        put_text_centered(frame, hint, fw // 2, fh - 72,
                          FONT_SMALL, 1.0, (70, 70, 70), 1)

    def _draw_empty_state(self, frame: np.ndarray, fw: int, fh: int):
        cy = fh // 2
        put_text_centered(frame, "Papka slaydov pusta", fw//2, cy-50,
                          FONT_BOLD, 1.0, COLOR_ACCENT, 2)
        put_text_centered(frame, f"Put': {SLIDES_DIR}", fw//2, cy,
                          FONT_SMALL, 1.0, COLOR_WHITE, 1)
        put_text_centered(frame, "Podderzhivayetsya: JPG, PNG, PDF, PPTX",
                          fw//2, cy+40, FONT_SMALL, 1.1, COLOR_GRAY, 1)

        # Инструкция по PDF
        if not _FITZ_OK:
            put_text_centered(frame, "dlya PDF: pip install pymupdf",
                              fw//2, cy+70, FONT_SMALL, 1.1, (0,140,255), 1)
        if not _PPTX_OK:
            put_text_centered(frame, "dlya PPTX: pip install python-pptx",
                              fw//2, cy+90, FONT_SMALL, 1.1, (0,140,255), 1)
